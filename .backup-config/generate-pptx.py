#!/usr/bin/env python3
"""
Node Governance Demo - PowerPoint Generator

Prerequisites:
    pip install python-pptx Pillow

Usage:
    cd .backup-config
    python generate-pptx.py

Output:
    node-governance-demo.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor  # Note: RGBColor not RgbColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# Configuration
SCREENSHOT_DIR = "demo-screenshots"
OUTPUT_FILE = "node-governance-demo.pptx"

# Colors
DARK_BG = RGBColor(26, 26, 46)       # #1a1a2e
ACCENT_RED = RGBColor(255, 107, 107)  # #ff6b6b
ACCENT_TEAL = RGBColor(78, 205, 196)  # #4ecdc4
ACCENT_YELLOW = RGBColor(254, 202, 87) # #feca57
GRAY = RGBColor(149, 165, 166)        # #95a5a6
WHITE = RGBColor(255, 255, 255)

def add_title_slide(prs, title, subtitle):
    """Add a title slide."""
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Background
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = DARK_BG
    background.line.fill.background()
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = ACCENT_RED
    p.alignment = PP_ALIGN.CENTER
    
    # Subtitle
    sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.5), Inches(9), Inches(0.8))
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(28)
    p.font.color.rgb = ACCENT_TEAL
    p.alignment = PP_ALIGN.CENTER
    
    return slide

def add_content_slide(prs, title, subtitle, bullets, image_path=None):
    """Add a content slide with optional image."""
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Background
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = DARK_BG
    background.line.fill.background()
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.7))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = ACCENT_TEAL
    
    # Subtitle
    sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.9), Inches(9), Inches(0.5))
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(18)
    p.font.color.rgb = GRAY
    
    # Bullets
    if image_path and os.path.exists(image_path):
        # Side-by-side layout
        bullet_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(4.2), Inches(4))
    else:
        # Full width
        bullet_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(4))
    
    tf = bullet_box.text_frame
    tf.word_wrap = True
    
    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"→ {bullet}"
        p.font.size = Pt(18)
        p.font.color.rgb = WHITE
        p.space_after = Pt(12)
    
    # Image
    if image_path and os.path.exists(image_path):
        try:
            slide.shapes.add_picture(
                image_path,
                Inches(4.8), Inches(1.5),
                width=Inches(5)
            )
        except Exception as e:
            print(f"Warning: Could not add image {image_path}: {e}")
    
    return slide

def add_summary_slide(prs, title, features):
    """Add a summary slide with feature grid."""
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Background
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = DARK_BG
    background.line.fill.background()
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = ACCENT_RED
    p.alignment = PP_ALIGN.CENTER
    
    # Feature grid (2x3)
    cols = 3
    box_width = 2.9
    box_height = 1.8
    start_x = 0.5
    start_y = 1.8
    gap = 0.2
    
    for i, (feat_title, feat_desc) in enumerate(features):
        row = i // cols
        col = i % cols
        
        x = start_x + col * (box_width + gap)
        y = start_y + row * (box_height + gap)
        
        # Box
        box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(x), Inches(y),
            Inches(box_width), Inches(box_height)
        )
        box.fill.solid()
        box.fill.fore_color.rgb = RGBColor(40, 40, 60)
        box.line.color.rgb = RGBColor(60, 60, 80)
        
        # Title
        title_tb = slide.shapes.add_textbox(
            Inches(x + 0.1), Inches(y + 0.2),
            Inches(box_width - 0.2), Inches(0.4)
        )
        tf = title_tb.text_frame
        p = tf.paragraphs[0]
        p.text = feat_title
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = ACCENT_TEAL
        p.alignment = PP_ALIGN.CENTER
        
        # Description
        desc_tb = slide.shapes.add_textbox(
            Inches(x + 0.1), Inches(y + 0.6),
            Inches(box_width - 0.2), Inches(1)
        )
        tf = desc_tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = feat_desc
        p.font.size = Pt(12)
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER
    
    # Questions
    q_box = slide.shapes.add_textbox(Inches(0.5), Inches(5.2), Inches(9), Inches(0.5))
    tf = q_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Questions?"
    p.font.size = Pt(24)
    p.font.color.rgb = GRAY
    p.alignment = PP_ALIGN.CENTER
    
    return slide


def main():
    print("Creating Node Governance Demo PowerPoint...")
    
    # Create presentation (16:9 aspect ratio)
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)
    
    # Slide 1: Title
    add_title_slide(
        prs,
        "Node Governance",
        "Control Which Nodes Can Be Used in Workflows"
    )
    print("  Added: Title slide")
    
    # Slide 2: Overview
    add_content_slide(
        prs,
        "Settings Overview",
        "Settings → Node Governance",
        [
            "Enable/Disable feature globally",
            "Priority notice explaining policy precedence",
            "Three tabs: Policies, Categories, Requests",
            "Priority: Project ALLOW > Global BLOCK > Global ALLOW > Project BLOCK"
        ],
        f"{SCREENSHOT_DIR}/demo-01-node-governance-overview.png"
    )
    print("  Added: Settings Overview")
    
    # Slide 3: Categories
    add_content_slide(
        prs,
        "Categories Tab",
        "Group related nodes for bulk policy management",
        [
            "AWS - 25 nodes",
            "Azure/Microsoft - 27 nodes",
            "Blocked External - 430 nodes",
            "Google - 55 nodes",
            "Block entire category instead of individual nodes"
        ],
        f"{SCREENSHOT_DIR}/demo-03-categories-tab.png"
    )
    print("  Added: Categories Tab")
    
    # Slide 4: Manage Nodes
    add_content_slide(
        prs,
        "Manage Category Nodes",
        "Add or remove nodes from a category",
        [
            "Search to find nodes to add",
            "View all assigned nodes",
            "Click X to remove a node",
            "Node type names shown for clarity"
        ],
        f"{SCREENSHOT_DIR}/demo-04-manage-nodes-modal.png"
    )
    print("  Added: Manage Nodes")
    
    # Slide 5: Policies
    add_content_slide(
        prs,
        "Policies Tab",
        "Define what's allowed and what's blocked",
        [
            "Status: Blocked (red) or Allowed (green)",
            "Target: Single Node or Category",
            "Scope: Global or Project-specific",
            "Edit or delete existing policies"
        ],
        f"{SCREENSHOT_DIR}/demo-05-policies-tab.png"
    )
    print("  Added: Policies Tab")
    
    # Slide 6: Create Policy
    add_content_slide(
        prs,
        "Create Policy",
        "Granular control over node access",
        [
            "Policy Type: Block or Allow",
            "Scope: Global or Specific Projects",
            "Target Type: Single Node or Category",
            "Target Value: e.g., n8n-nodes-base.ssh"
        ],
        f"{SCREENSHOT_DIR}/demo-06-create-policy-modal.png"
    )
    print("  Added: Create Policy")
    
    # Slide 7: Node Creator Integration
    add_content_slide(
        prs,
        "Node Creator Integration",
        "Real-time governance status when adding nodes",
        [
            "Blocked nodes show 'Blocked' badge",
            "'Request access' link appears",
            "Users can still see blocked nodes",
            "Prevents accidental use of restricted nodes"
        ],
        f"{SCREENSHOT_DIR}/demo-09-gmail-blocked.png"
    )
    print("  Added: Node Creator Integration")
    
    # Slide 8: Request Access
    add_content_slide(
        prs,
        "Request Node Access",
        "Users provide business justification",
        [
            "Project - Which project needs access",
            "Workflow Name - Optional context",
            "Justification - Required business reason",
            "Creates audit trail for compliance"
        ],
        f"{SCREENSHOT_DIR}/demo-10-request-access-modal.png"
    )
    print("  Added: Request Access")
    
    # Slide 9: Admin Review
    add_content_slide(
        prs,
        "Admin Review",
        "Approve or reject access requests",
        [
            "Pending requests appear with badge count",
            "View user info and justification",
            "Approve or Reject with comment",
            "Full audit trail maintained"
        ],
        f"{SCREENSHOT_DIR}/demo-11-requests-tab.png"
    )
    print("  Added: Admin Review")
    
    # Slide 10: Summary
    add_summary_slide(
        prs,
        "Key Benefits",
        [
            ("Security", "Block risky nodes like SSH, Execute Command, Code execution"),
            ("Compliance", "Audit trail for all access requests and approvals"),
            ("Flexibility", "Global or project-specific policies with priority system"),
            ("Visibility", "Users see blocked status in real-time when adding nodes"),
            ("Categories", "Manage hundreds of nodes with bulk category policies"),
            ("Self-Service", "Users can request access with business justification"),
        ]
    )
    print("  Added: Summary")
    
    # Save
    prs.save(OUTPUT_FILE)
    print(f"\nSuccess! Created: {OUTPUT_FILE}")
    print(f"Location: {os.path.abspath(OUTPUT_FILE)}")


if __name__ == "__main__":
    main()
